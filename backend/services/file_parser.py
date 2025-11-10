import os
import asyncio
import logging
import hashlib
import shutil
import fnmatch
import datetime as dt
from concurrent.futures import ThreadPoolExecutor
from typing import Dict, Any, List, Optional, Callable

# --- Dependency Check ---
try:
    import aiofiles
    import magic
    from diskcache import Cache
    from tqdm.asyncio import tqdm_asyncio
except ImportError as e:
    logging.basicConfig(level=logging.CRITICAL, format='%(asctime)s - %(levelname)s - %(message)s')
    logging.critical(f"Missing critical dependency: {e}. Please run 'pip install -r requirements.txt'")
    exit(1)

# --- Constants ---
CACHE_DIR = ".file_parser_cache"

class FileParser:
    """
    A unified, asynchronous parser for a wide range of documents and images, featuring
    persistent caching, MIME type detection, recursive directory parsing, configurable OCR,
    and seamless NLP integration.
    """
    _cache = Cache(CACHE_DIR)
    _cpu_executor = ThreadPoolExecutor(max_workers=min(32, (os.cpu_count() or 1) + 4))

    # --- Core Class Methods (unchanged) ---
    @staticmethod
    def set_log_level(level: str = "INFO"):
        numeric_level = getattr(logging, level.upper(), None)
        if not isinstance(numeric_level, int):
            raise ValueError(f"Invalid log level: {level}")
        logging.basicConfig(level=numeric_level, format='%(asctime)s - %(levelname)s - %(message)s', force=True)
        logging.info(f"Logging level set to {level.upper()}")

    @staticmethod
    def clear_cache():
        logging.info(f"Clearing cache at {CACHE_DIR}")
        FileParser._cache.clear()
        logging.info("Cache cleared.")

    @staticmethod
    async def prune_cache(valid_filepaths: List[str]):
        logging.info("Starting cache pruning...")
        hash_tasks = [FileParser._get_file_hash(fp) for fp in valid_filepaths if os.path.isfile(fp)]
        valid_hashes = set(await asyncio.gather(*hash_tasks))
        cached_keys = set(FileParser._cache.iterkeys())
        keys_to_delete = cached_keys - valid_hashes
        if not keys_to_delete:
            logging.info("Cache is clean. No keys to prune.")
            return
        logging.warning(f"Pruning {len(keys_to_delete)} stale entries from cache.")
        for key in keys_to_delete:
            del FileParser._cache[key]
        logging.info("Cache pruning complete.")

    @staticmethod
    async def _get_file_hash(filepath: str) -> str:
        hasher = hashlib.sha256()
        try:
            async with aiofiles.open(filepath, 'rb') as f:
                while chunk := await f.read(8192):
                    hasher.update(chunk)
        except FileNotFoundError: return ""
        return hasher.hexdigest()

    @staticmethod
    def _get_mime_type(filepath: str) -> str:
        try:
            mime = magic.Magic(mime=True)
            return mime.from_file(filepath)
        except Exception: return ""

    @staticmethod
    async def _run_sync_in_executor(func, *args):
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(FileParser._cpu_executor, func, *args)

    # --- Individual Parsers with Enhanced Logging ---

    @staticmethod
    async def _parse_txt(filepath: str, warnings: List[str], **kwargs) -> str:
        try: import chardet
        except ImportError:
            warnings.append("`chardet` not installed. Defaulting to utf-8.")
            logging.debug(f"Parsing TXT '{os.path.basename(filepath)}' with default utf-8 encoding.")
            async with aiofiles.open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                return await f.read()
        
        async with aiofiles.open(filepath, 'rb') as f: raw_data = await f.read()
        detected = chardet.detect(raw_data)
        encoding = detected['encoding'] or 'utf-8'
        logging.debug(f"Parsing TXT '{os.path.basename(filepath)}' with detected encoding: {encoding}.")
        return raw_data.decode(encoding, errors='replace')

    @staticmethod
    async def _parse_docx(filepath: str, warnings: List[str], **kwargs) -> str:
        def docx_task():
            from docx import Document
            doc = Document(filepath)
            paragraphs = [p.text for p in doc.paragraphs if p.text]
            logging.debug(f"Extracted {len(paragraphs)} paragraphs from DOCX '{os.path.basename(filepath)}'.")
            return "\n".join(paragraphs)
        return await FileParser._run_sync_in_executor(docx_task)
    
    @staticmethod
    async def _parse_xlsx(filepath: str, warnings: List[str], **kwargs) -> str:
        """Parses .xlsx files using openpyxl."""
        def xlsx_task():
            try: from openpyxl import load_workbook
            except ImportError:
                warnings.append("`openpyxl` not installed. Cannot parse .xlsx files.")
                return ""
            workbook = load_workbook(filepath, read_only=True, data_only=True)
            all_text = []
            for sheet in workbook:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            all_text.append(str(cell.value).strip())
            logging.debug(f"Extracted text from {len(workbook.sheetnames)} sheets in XLSX '{os.path.basename(filepath)}'.")
            return "\n".join(filter(None, all_text))
        return await FileParser._run_sync_in_executor(xlsx_task)

    @staticmethod
    async def _parse_pptx(filepath: str, warnings: List[str], **kwargs) -> str:
        def pptx_task():
            from pptx import Presentation
            pres = Presentation(filepath)
            text_runs = [p.text for s in pres.slides for sh in s.shapes if sh.has_text_frame for p in sh.text_frame.paragraphs]
            logging.debug(f"Extracted text from {len(pres.slides)} slides in PPTX '{os.path.basename(filepath)}'.")
            return "\n".join(text_runs)
        return await FileParser._run_sync_in_executor(pptx_task)

    @staticmethod
    async def _parse_odt(filepath: str, warnings: List[str], **kwargs) -> str:
        def odt_task():
            from odf import text, teletype
            from odf.opendocument import load
            doc = load(filepath)
            paragraphs = doc.getElementsByType(text.P)
            logging.debug(f"Extracted {len(paragraphs)} paragraphs from ODT '{os.path.basename(filepath)}'.")
            return "\n".join(teletype.extractText(p) for p in paragraphs)
        return await FileParser._run_sync_in_executor(odt_task)

    @staticmethod
    async def _parse_epub(filepath: str, warnings: List[str], **kwargs) -> str:
        def epub_task():
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            book = epub.read_epub(filepath)
            items = list(book.get_items_of_type(ebooklib.ITEM_DOCUMENT))
            logging.debug(f"Extracted {len(items)} sections from EPUB '{os.path.basename(filepath)}'.")
            return "\n".join(BeautifulSoup(i.get_body_content(), 'html.parser').get_text() for i in items)
        return await FileParser._run_sync_in_executor(epub_task)

    @staticmethod
    async def _parse_image(filepath: str, warnings: List[str], **ocr_kwargs) -> str:
        def image_ocr_task():
            try:
                from PIL import Image
                import pytesseract
                logging.debug(f"Performing OCR on image '{os.path.basename(filepath)}' with config: {ocr_kwargs}")
                return pytesseract.image_to_string(Image.open(filepath), **ocr_kwargs)
            except Exception as e:
                warnings.append(f"Image OCR failed: {e}")
                return ""
        return await FileParser._run_sync_in_executor(image_ocr_task)

    @staticmethod
    async def _parse_pdf(filepath: str, warnings: List[str], **ocr_kwargs) -> str:
        def pdf_task():
            import pdfplumber
            with pdfplumber.open(filepath) as pdf:
                logging.debug(f"Parsing {len(pdf.pages)} pages from PDF '{os.path.basename(filepath)}'.")
                return "\n".join(page.extract_text() or "" for page in pdf.pages)
        full_text = await FileParser._run_sync_in_executor(pdf_task)
        if len(full_text.strip()) < 100:
            msg = f"Low text from PDF '{os.path.basename(filepath)}'; attempting OCR."
            warnings.append(msg)
            logging.warning(msg)
            def ocr_pdf_task():
                from pdf2image import convert_from_path
                import pytesseract
                logging.debug(f"Performing PDF-OCR on '{os.path.basename(filepath)}' with config: {ocr_kwargs}")
                pages = convert_from_path(filepath, dpi=ocr_kwargs.get('dpi', 200))
                return "\n".join([pytesseract.image_to_string(p, **ocr_kwargs) for p in pages])
            return await FileParser._run_sync_in_executor(ocr_pdf_task)
        return full_text
        
    @staticmethod
    async def _parse_with_fallback(filepath: str, warnings: List[str], **kwargs) -> str:
        logging.debug(f"Using fallback parser for '{os.path.basename(filepath)}'.")
        def fallback_task():
            import textract
            return textract.process(filepath).decode('utf-8', errors='replace')
        return await FileParser._run_sync_in_executor(fallback_task)


    @staticmethod
    async def parse_file(filepath: str, use_cache: bool = True, ocr_config: Optional[Dict] = None) -> Dict[str, Any]:
        if not os.path.isfile(filepath): raise FileNotFoundError(f"Not a valid file: {filepath}")
        file_hash = await FileParser._get_file_hash(filepath)
        if use_cache and file_hash in FileParser._cache: return FileParser._cache[file_hash]

        mime_type = FileParser._get_mime_type(filepath)
        _, extension = os.path.splitext(filepath)
        extension = extension.lower()
        ocr_config = ocr_config or {}

        parser_map = {
            'application/pdf': (FileParser._parse_pdf, ocr_config.get('pdf', {})),
            'image/png': (FileParser._parse_image, ocr_config.get('image', {})),
            'image/jpeg': (FileParser._parse_image, ocr_config.get('image', {})),
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': (FileParser._parse_docx, {}),
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': (FileParser._parse_pptx, {}),
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': (FileParser._parse_xlsx, {}),
            'application/vnd.oasis.opendocument.text': (FileParser._parse_odt, {}),
            'application/epub+zip': (FileParser._parse_epub, {}),
            'application/msword': (FileParser._parse_with_fallback, {}),
            'application/vnd.ms-powerpoint': (FileParser._parse_with_fallback, {}),
            'application/rtf': (FileParser._parse_with_fallback, {}), 'text/rtf': (FileParser._parse_with_fallback, {}),
            'text/plain': (FileParser._parse_txt, {}),
        }
        ext_map = {
            '.docx': (FileParser._parse_docx, {}), '.pptx': (FileParser._parse_pptx, {}), '.xlsx': (FileParser._parse_xlsx, {}),
            '.odt': (FileParser._parse_odt, {}), '.epub': (FileParser._parse_epub, {}),
            '.doc': (FileParser._parse_with_fallback, {}), '.ppt': (FileParser._parse_with_fallback, {}),
            '.txt': (FileParser._parse_txt, {}),
        }

        parser_func, kwargs = parser_map.get(mime_type, ext_map.get(extension, (FileParser._parse_with_fallback, {})))
        parsing_warnings, extracted_text = [], ""
        logging.info(f"Parsing '{os.path.basename(filepath)}' using '{parser_func.__name__}'")
        try: extracted_text = await parser_func(filepath, parsing_warnings, **kwargs)
        except Exception as e:
            msg = f"Unexpected error during '{parser_func.__name__}': {e}"
            logging.error(msg, exc_info=True)
            parsing_warnings.append(msg)
        result = { "filepath": filepath, "file_hash": file_hash, "text": extracted_text, "extraction_method": parser_func.__name__, "parsing_warnings": parsing_warnings }
        if use_cache: FileParser._cache[file_hash] = result
        return result

    @staticmethod
    async def bulk_parse_files(paths: List[str], patterns: List[str] = ['*.*'], nlp_function: Optional[Callable[[str], Dict]] = None, ocr_config: Optional[Dict] = None) -> List[Dict[str, Any]]:
        # Stage 1: File Discovery
        filepaths = []
        for path in paths:
            if os.path.isfile(path): filepaths.append(path)
            elif os.path.isdir(path):
                for root, _, files in os.walk(path):
                    for file in files:
                        if any(fnmatch.fnmatch(file, p) for p in patterns):
                            filepaths.append(os.path.join(root, file))
        
        # Stage 2: Parallel File Parsing
        parsing_tasks = [FileParser.parse_file(fp, ocr_config=ocr_config) for fp in filepaths]
        parsed_results = []
        for future in tqdm_asyncio.as_completed(parsing_tasks, total=len(parsing_tasks), desc="Parsing files"):
            try: parsed_results.append(await future)
            except Exception as e: logging.error(f"A file parsing task failed: {e}", exc_info=True)
        
        # Stage 3: Parallel NLP Enrichment (if applicable)
        if not nlp_function: return parsed_results

        nlp_targets = [(i, res) for i, res in enumerate(parsed_results) if res.get('text')]
        if not nlp_targets: return parsed_results
        
        async def run_nlp_safe(text: str, filepath: str):
            try: return await FileParser._run_sync_in_executor(nlp_function, text)
            except Exception as e:
                msg = f"NLP function failed for {os.path.basename(filepath)}: {e}"
                logging.error(msg, exc_info=True)
                return {"nlp_error": msg}

        nlp_tasks = [run_nlp_safe(res['text'], res['filepath']) for _, res in nlp_targets]
        logging.info(f"Starting parallel NLP enrichment for {len(nlp_tasks)} documents.")
        
        nlp_results_iter = tqdm_asyncio.as_completed(nlp_tasks, total=len(nlp_tasks), desc="Enriching text (NLP)")
        nlp_data_list = [await future for future in nlp_results_iter]

        for (original_index, _), nlp_data in zip(nlp_targets, nlp_data_list):
            if "nlp_error" in nlp_data:
                parsed_results[original_index]['parsing_warnings'].append(nlp_data["nlp_error"])
            else:
                parsed_results[original_index].update(nlp_data)
        
        return parsed_results

async def main():
    def analyze_text_mock(text: str) -> dict:
        return {"summary": text[:30].strip() + "...", "word_count": len(text.split())}

    FileParser.set_log_level("DEBUG") # Use DEBUG to see detailed parser logs
    test_dir = "test_unified_parser"
    if os.path.exists(test_dir): shutil.rmtree(test_dir)
    os.makedirs(test_dir)
    
    dummy_files = {
        "report.txt": "Text content.", "invoice.png": b"png", "manual.pdf": b"pdf",
        "proposal.docx": b"docx", "meeting.pptx": b"pptx", "data.xlsx": b"xlsx",
        "draft.odt": b"odt", "ebook.epub": b"epub", "legacy.doc": b"doc", 
        "archive.ppt": b"ppt", "notes.rtf": b"rtf"
    }
    for name, content in dummy_files.items():
        with open(os.path.join(test_dir, name), "wb" if isinstance(content, bytes) else "w") as f:
            f.write(content)

    print("\n--- Starting Unified Bulk Processing with NLP ---")
    processed_data = await FileParser.bulk_parse_files(paths=[test_dir], nlp_function=analyze_text_mock)
    print(f"\n--- Finished processing {len(processed_data)} files ---\n")
    for result in sorted(processed_data, key=lambda x: x['filepath']):
        print(f"  - Parsed: {os.path.basename(result['filepath']):<15} -> Method: {result['extraction_method']}")

    shutil.rmtree(test_dir)
    FileParser._cache.close()

if __name__ == '__main__':
    if not os.path.exists(CACHE_DIR): os.makedirs(CACHE_DIR)
    asyncio.run(main())

